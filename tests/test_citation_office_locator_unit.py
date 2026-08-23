from __future__ import annotations

import asyncio
from io import BytesIO
from types import SimpleNamespace

from library.agent import runtime


def _file(**overrides):
    values = {
        "id": "file-1",
        "storage_key": "object-key",
        "sha256": "a" * 64,
        "mime_type": "application/octet-stream",
        "original_ext": "bin",
        "kind": None,
    }
    values.update(overrides)
    return SimpleNamespace(**values)


def test_citation_query_string_combines_native_anchor_and_quote() -> None:
    pdf = _file(mime_type="application/pdf", original_ext="pdf", kind="text")
    docx = _file(
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        original_ext="docx",
        kind="docx",
    )
    pptx = _file(
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        original_ext="pptx",
        kind="text",
    )
    xlsx = _file(
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        original_ext="xlsx",
        kind="table",
    )

    assert runtime._pick_query_string(
        pdf,
        "quick final tail",
        "3",
        located_pdf_page=7,
    ) == "?page=7&q=quick+final+tail"
    assert runtime._pick_query_string(
        docx,
        "contract clause",
        None,
        located_docx_block=42,
    ) == "?block=42&q=contract+clause"
    assert runtime._pick_query_string(
        pptx,
        "slide bullet",
        "2",
        located_pptx_slide=6,
    ) == "?page=6&q=slide+bullet"
    assert runtime._pick_query_string(
        xlsx,
        "budget variance",
        None,
        located_xlsx_cell=runtime._SpreadsheetLocator(
            sheet="FY 2026",
            cell="B14",
            row=14,
        ),
    ) == "?sheet=FY+2026&cell=B14&q=budget+variance"
    assert runtime._pick_query_string(
        xlsx,
        "split across cells",
        None,
        located_xlsx_cell=runtime._SpreadsheetLocator(sheet="FY 2026", row=9),
    ) == "?sheet=FY+2026&row=9&q=split+across+cells"


def test_citation_locators_find_office_anchors(monkeypatch) -> None:
    bodies: dict[str, bytes] = {}

    async def fake_body(file):
        return bodies[file.storage_key]

    monkeypatch.setattr(runtime, "_read_file_body_for_locator", fake_body)

    from docx import Document as DocxDocument

    document = DocxDocument()
    document.add_paragraph("Intro")
    document.add_paragraph("The contract clause appears here.")
    docx_out = BytesIO()
    document.save(docx_out)
    bodies["docx-key"] = docx_out.getvalue()

    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    slide_one = deck.slides.add_slide(deck.slide_layouts[5])
    slide_one.shapes.title.text = "Roadmap"
    slide_one.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(5), Inches(1),
    ).text = "Alpha"
    slide_two = deck.slides.add_slide(deck.slide_layouts[5])
    slide_two.shapes.title.text = "Budget"
    slide_two.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(5), Inches(1),
    ).text = "Slide bullet target"
    pptx_out = BytesIO()
    deck.save(pptx_out)
    bodies["pptx-key"] = pptx_out.getvalue()

    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "FY 2026"
    sheet["B14"] = "budget variance target"
    sheet["A9"] = "split"
    sheet["B9"] = "across cells"
    xlsx_out = BytesIO()
    workbook.save(xlsx_out)
    workbook.close()
    bodies["xlsx-key"] = xlsx_out.getvalue()

    docx_file = _file(storage_key="docx-key", original_ext="docx")
    pptx_file = _file(storage_key="pptx-key", original_ext="pptx")
    xlsx_file = _file(storage_key="xlsx-key", original_ext="xlsx")

    assert asyncio.run(
        runtime._locate_docx_quote_block(docx_file, "contract clause"),
    ) == 2
    assert asyncio.run(
        runtime._locate_pptx_quote_slide(pptx_file, "slide bullet target"),
    ) == 2
    assert asyncio.run(
        runtime._locate_xlsx_quote_cell(xlsx_file, "budget variance"),
    ) == runtime._SpreadsheetLocator(sheet="FY 2026", cell="B14", row=14)
    assert asyncio.run(
        runtime._locate_xlsx_quote_cell(xlsx_file, "split across cells"),
    ) == runtime._SpreadsheetLocator(sheet="FY 2026", row=9)
