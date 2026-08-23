from __future__ import annotations

import io
import zipfile
from types import SimpleNamespace

import pytest

from library.pipelines.archive import ArchivePipeline
from library.pipelines.docx import DocxPipeline
from library.pipelines.markitdown import MarkItDownPipeline
from library.pipelines.pptx import PptxPipeline
from library.pipelines.spreadsheet import SpreadsheetPipeline


class _MemoryStorage:
    def __init__(self, body: bytes) -> None:
        self.body = body

    async def get(self, key: str):
        assert key == "file-key"
        yield self.body


def _xlsx_with_late_needle() -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "data"
    ws.append(["row", "value"])
    for idx in range(1, 260):
        ws.append([idx, f"value-{idx}"])
    ws.append([260, "needle-after-row-cap"])
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def _docx_with_heading() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("Contract Heading", level=1)
    doc.add_paragraph("Heading body line")
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _docx_with_table() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Product definition")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Indicator"
    table.cell(0, 1).text = "Strategy"
    table.cell(1, 0).text = "Comfort"
    table.cell(1, 1).text = "L"
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _pptx_with_heading() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Contract Slide"
    box = slide.shapes.add_textbox(914400, 1371600, 5486400, 914400)
    box.text = "Slide body line"
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _pptx_with_empty_tail() -> bytes:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(_pptx_with_heading()))
    prs.slides.add_slide(prs.slide_layouts[6])
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


@pytest.mark.asyncio
async def test_markitdown_heading_read_uses_full_extracted_text(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import library.pipelines.markitdown as mod

    monkeypatch.setattr(
        mod,
        "_convert_bytes_with_markitdown",
        lambda body, suffix: "A" * 125_000 + "\n# Late Chapter\nneedle-after-index-cap",
    )

    result = await MarkItDownPipeline().read_segment_from_bytes(
        b"source bytes",
        {"heading": "Late Chapter", "max_chars": 200},
        filename="book.epub",
    )

    assert result.error is None
    assert "needle-after-index-cap" in result.text
    assert result.extras["located_via"] == "body-heading-scan"


@pytest.mark.asyncio
async def test_spreadsheet_read_uses_full_rows_not_ingest_sample() -> None:
    result = await SpreadsheetPipeline().read_segment(
        file_row=SimpleNamespace(storage_key="file-key"),
        args={"pattern": "needle-after-row-cap", "context_lines": 0},
        storage=_MemoryStorage(_xlsx_with_late_needle()),
    )

    assert result.error is None
    assert "needle-after-row-cap" in result.text
    assert result.extras["total_matches"] == 1


@pytest.mark.asyncio
async def test_docx_heading_and_line_reads_use_complete_extracted_text() -> None:
    pipeline = DocxPipeline()
    row = SimpleNamespace(storage_key="file-key", description=None)

    heading = await pipeline.read_segment(
        file_row=row,
        args={"heading": "Contract Heading", "max_chars": 200},
        storage=_MemoryStorage(_docx_with_heading()),
    )
    assert heading.error is None
    assert "Heading body line" in heading.text

    line = await pipeline.read_segment(
        file_row=row,
        args={"line_start": 1, "line_end": 2},
        storage=_MemoryStorage(_docx_with_heading()),
    )
    assert line.error is None
    assert "Contract Heading" in line.text
    assert line.extras["line_start"] == 1


@pytest.mark.asyncio
async def test_pptx_heading_and_line_reads_use_complete_extracted_text() -> None:
    pipeline = PptxPipeline()
    row = SimpleNamespace(storage_key="file-key", description=None)

    heading = await pipeline.read_segment(
        file_row=row,
        args={"heading": "Slide 1: Contract Slide", "max_chars": 300},
        storage=_MemoryStorage(_pptx_with_heading()),
    )
    assert heading.error is None
    assert "Slide body line" in heading.text

    line = await pipeline.read_segment(
        file_row=row,
        args={"line_start": 1, "line_end": 2},
        storage=_MemoryStorage(_pptx_with_heading()),
    )
    assert line.error is None
    assert "Contract Slide" in line.text
    assert line.extras["line_start"] == 1


@pytest.mark.asyncio
async def test_docx_question_prefers_extracted_text(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    async def forbidden_vision(**_kwargs):
        raise AssertionError("readable DOCX text must bypass vision")

    monkeypatch.setattr(
        "library.pipelines.docx.answer_document_image_question",
        forbidden_vision,
    )
    result = await DocxPipeline().read_segment(
        file_row=SimpleNamespace(storage_key="file-key", description=None),
        args={"question": "What is the strategy?"},
        storage=_MemoryStorage(_docx_with_table()),
    )

    assert result.error is None
    assert result.extras["mode"] == "docx_text_question"
    assert result.extras["answered_by"] == "docx_extracted_text"
    assert result.extras["source_text_preserved"] is True
    assert "Product definition" in result.text
    assert "Comfort | L" in result.text


@pytest.mark.asyncio
async def test_pptx_question_prefers_extracted_text(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    async def forbidden_vision(**_kwargs):
        raise AssertionError("readable PPTX text must bypass vision")

    monkeypatch.setattr(
        "library.pipelines.pptx.answer_document_image_question",
        forbidden_vision,
    )
    result = await PptxPipeline().read_segment(
        file_row=SimpleNamespace(storage_key="file-key", description=None),
        args={"question": "What does the slide say?"},
        storage=_MemoryStorage(_pptx_with_heading()),
    )

    assert result.error is None
    assert result.extras["mode"] == "pptx_text_question"
    assert result.extras["answered_by"] == "pptx_extracted_text"
    assert result.extras["source_text_preserved"] is True
    assert "Slide body line" in result.text


@pytest.mark.asyncio
async def test_pptx_empty_tail_does_not_replace_earlier_source_text() -> None:
    result = await PptxPipeline().read_segment(
        file_row=SimpleNamespace(
            storage_key="file-key",
            description={
                "document_vision": {
                    "text": "Persisted image-only evidence.",
                    "images": [],
                }
            },
        ),
        args={"slide_start": 1, "slide_end": 2},
        storage=_MemoryStorage(_pptx_with_empty_tail()),
    )

    assert result.error is None
    assert "Contract Slide" in result.text
    assert "Slide body line" in result.text
    assert "# Slide 2\n(no extractable text)" in result.text


@pytest.mark.asyncio
async def test_archive_member_read_dispatches_to_full_inner_reader() -> None:
    archive_body = io.BytesIO()
    with zipfile.ZipFile(archive_body, "w") as zf:
        zf.writestr("reports/data.xlsx", _xlsx_with_late_needle())

    result = await ArchivePipeline().read_segment(
        file_row=SimpleNamespace(storage_key="file-key", original_ext=".zip", id="file-id"),
        args={
            "member_path": "reports/data.xlsx",
            "pattern": "needle-after-row-cap",
            "context_lines": 0,
        },
        storage=_MemoryStorage(archive_body.getvalue()),
    )

    assert result.error is None
    assert "needle-after-row-cap" in result.text
    assert result.extras["total_matches"] == 1
